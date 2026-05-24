"""
ArchDefend — Export Engine
Generates PDF, PPTX, Markdown, and HTML reports from analysis data.
Uses WeasyPrint (PDF) and python-pptx (PPTX).
"""

import io
import json
import logging
import os
from datetime import datetime
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)


# ── HTML Template for PDF ─────────────────────────────────────────────────────

PDF_HTML_TEMPLATE = """
<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<style>
  @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700&family=JetBrains+Mono:wght@400;500&display=swap');

  :root {{
    --bg: #0A0A0A;
    --surface: #111111;
    --accent: #22D3EE;
    --text: #F4F4F5;
    --muted: #A1A1AA;
    --border: #1F1F23;
  }}

  * {{ box-sizing: border-box; margin: 0; padding: 0; }}

  body {{
    font-family: 'Inter', system-ui, sans-serif;
    background: #FFFFFF;
    color: #111111;
    font-size: 11pt;
    line-height: 1.6;
  }}

  .cover {{
    background: linear-gradient(135deg, #0A0A0A 0%, #111827 100%);
    color: #F4F4F5;
    padding: 60px 48px;
    min-height: 100vh;
    display: flex;
    flex-direction: column;
    justify-content: center;
  }}

  .cover-logo {{
    width: 48px;
    height: 48px;
    background: #22D3EE;
    border-radius: 10px;
    display: flex;
    align-items: center;
    justify-content: center;
    margin-bottom: 48px;
    font-size: 24px;
  }}

  .cover h1 {{
    font-size: 40pt;
    font-weight: 700;
    margin-bottom: 16px;
    letter-spacing: -1px;
  }}

  .cover h1 span {{ color: #22D3EE; }}
  .cover p {{ color: #A1A1AA; font-size: 13pt; margin-bottom: 8px; }}

  .cover .meta {{
    margin-top: 48px;
    padding-top: 24px;
    border-top: 1px solid #1F1F23;
    color: #6B7280;
    font-size: 9pt;
    font-family: 'JetBrains Mono', monospace;
  }}

  .page {{
    padding: 48px;
    page-break-after: always;
  }}

  h2 {{
    font-size: 22pt;
    font-weight: 700;
    margin-bottom: 24px;
    padding-bottom: 12px;
    border-bottom: 2px solid #22D3EE;
    color: #111111;
  }}

  h3 {{
    font-size: 14pt;
    font-weight: 600;
    margin: 24px 0 12px;
    color: #111111;
  }}

  p {{ margin-bottom: 12px; color: #374151; }}

  .score-card {{
    display: inline-block;
    padding: 24px 32px;
    border-radius: 12px;
    border: 2px solid #E5E7EB;
    text-align: center;
    margin: 8px;
  }}

  .score-number {{
    font-size: 36pt;
    font-weight: 700;
    color: #22D3EE;
  }}

  .score-label {{
    font-size: 9pt;
    color: #6B7280;
    margin-top: 4px;
  }}

  .badge {{
    display: inline-block;
    padding: 3px 10px;
    border-radius: 20px;
    font-size: 8pt;
    font-weight: 600;
    margin-right: 6px;
  }}

  .badge-critical {{ background: #FEE2E2; color: #DC2626; }}
  .badge-high {{ background: #FEF3C7; color: #D97706; }}
  .badge-medium {{ background: #FEF9C3; color: #CA8A04; }}
  .badge-low {{ background: #D1FAE5; color: #059669; }}
  .badge-info {{ background: #DBEAFE; color: #2563EB; }}

  .finding {{
    border: 1px solid #E5E7EB;
    border-radius: 8px;
    padding: 16px;
    margin-bottom: 12px;
    background: #F9FAFB;
  }}

  .finding-title {{
    font-weight: 600;
    font-size: 10pt;
    margin-bottom: 6px;
  }}

  .finding-meta {{
    font-size: 8pt;
    color: #6B7280;
    font-family: 'JetBrains Mono', monospace;
    margin-bottom: 8px;
  }}

  .code-block {{
    background: #0A0A0A;
    color: #22D3EE;
    padding: 16px;
    border-radius: 8px;
    font-family: 'JetBrains Mono', monospace;
    font-size: 8pt;
    overflow: hidden;
    margin: 12px 0;
  }}

  .qa-item {{
    border: 1px solid #E5E7EB;
    border-radius: 8px;
    padding: 20px;
    margin-bottom: 16px;
  }}

  .qa-question {{
    font-weight: 600;
    font-size: 10pt;
    margin-bottom: 12px;
    color: #111111;
  }}

  .qa-answer {{
    color: #374151;
    font-size: 9pt;
    line-height: 1.7;
  }}

  .tech-pill {{
    display: inline-block;
    padding: 4px 12px;
    border-radius: 20px;
    border: 1px solid #E5E7EB;
    font-size: 8pt;
    color: #374151;
    margin: 3px;
    background: #F9FAFB;
  }}

  table {{
    width: 100%;
    border-collapse: collapse;
    margin: 16px 0;
  }}

  th {{
    background: #F9FAFB;
    padding: 10px 12px;
    text-align: left;
    font-size: 8pt;
    font-weight: 600;
    color: #6B7280;
    border-bottom: 1px solid #E5E7EB;
    text-transform: uppercase;
    letter-spacing: 0.5px;
  }}

  td {{
    padding: 10px 12px;
    font-size: 9pt;
    border-bottom: 1px solid #F3F4F6;
    color: #374151;
  }}

  .footer {{
    text-align: center;
    color: #9CA3AF;
    font-size: 8pt;
    padding: 24px;
    border-top: 1px solid #E5E7EB;
    font-family: 'JetBrains Mono', monospace;
  }}
</style>
</head>
<body>

<!-- Cover Page -->
<div class="cover">
  <div class="cover-logo">🛡</div>
  <h1>Arch<span>Defend</span></h1>
  <p>Codebase Intelligence Report</p>
  <p style="color: #22D3EE; font-size: 18pt; font-weight: 600; margin-top: 16px;">{repo_name}</p>
  <p>{repo_url}</p>
  <div class="meta">
    Generated: {generated_at}<br>
    By EMEMZYVISUALS DIGITALS · archdefend.io
  </div>
</div>

<!-- Executive Summary -->
<div class="page">
  <h2>Executive Summary</h2>

  <div style="margin-bottom: 32px;">
    <div class="score-card">
      <div class="score-number">{scalability_score}</div>
      <div class="score-label">Scalability Score</div>
    </div>
    <div class="score-card">
      <div class="score-number">{readiness_score}</div>
      <div class="score-label">Production Readiness</div>
    </div>
    <div class="score-card">
      <div class="score-number" style="color: {security_color};">{security_count}</div>
      <div class="score-label">Security Findings</div>
    </div>
    <div class="score-card">
      <div class="score-number">{file_count}</div>
      <div class="score-label">Files Analyzed</div>
    </div>
  </div>

  <h3>Architecture Summary</h3>
  <p>{architecture_summary}</p>

  <h3>Technology Stack</h3>
  {tech_stack_html}
</div>

<!-- Security Analysis -->
<div class="page">
  <h2>Security Analysis</h2>
  {security_html}
</div>

<!-- API Inventory -->
<div class="page">
  <h2>API Inventory</h2>
  {api_html}
</div>

<!-- Interview Defense Questions -->
<div class="page">
  <h2>Interview Defense Questions</h2>
  <p style="color: #6B7280; margin-bottom: 24px;">Architecture-specific Q&A grounded in this codebase.</p>
  {interview_html}
</div>

<!-- Recommendations -->
<div class="page">
  <h2>Recommendations</h2>
  {recommendations_html}
  <div class="footer">ArchDefend · archdefend.io · Generated by EMEMZYVISUALS DIGITALS</div>
</div>

</body>
</html>
"""


class ExportEngine:
    """Generates PDF, PPTX, Markdown, and HTML reports."""

    async def generate_pdf(self, report_data: dict, analysis_meta: dict) -> bytes:
        """Generate PDF report using WeasyPrint."""
        try:
            from weasyprint import HTML, CSS
        except ImportError:
            raise RuntimeError("WeasyPrint not installed. Run: pip install weasyprint")

        html_content = self._build_pdf_html(report_data, analysis_meta)
        pdf_bytes = HTML(string=html_content).write_pdf()
        logger.info(f"PDF generated: {len(pdf_bytes):,} bytes")
        return pdf_bytes

    async def generate_pptx(self, report_data: dict, analysis_meta: dict) -> bytes:
        """Generate PPTX presentation using python-pptx."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            raise RuntimeError("python-pptx not installed. Run: pip install python-pptx")

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Color constants
        BLACK = RGBColor(10, 10, 10)
        ACCENT = RGBColor(34, 211, 238)
        WHITE = RGBColor(244, 244, 245)
        MUTED = RGBColor(161, 161, 170)
        SURFACE = RGBColor(17, 17, 17)

        def add_dark_slide(title_text: str, layout_idx: int = 6) -> object:
            """Add a dark-themed slide."""
            slide_layout = prs.slide_layouts[layout_idx]
            slide = prs.slides.add_slide(slide_layout)
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = BLACK
            return slide

        def add_title(slide, text: str, y=Inches(0.5), size=Pt(36)):
            txBox = slide.shapes.add_textbox(Inches(0.8), y, Inches(11.5), Inches(1.2))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.bold = True
            run.font.size = size
            run.font.color.rgb = WHITE

        def add_body(slide, text: str, y=Inches(1.8), color=None):
            txBox = slide.shapes.add_textbox(Inches(0.8), y, Inches(11.5), Inches(4.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.size = Pt(13)
            run.font.color.rgb = color or MUTED

        # Slide 1: Cover
        slide = add_dark_slide("Cover")
        add_title(slide, "ArchDefend", y=Inches(1.5), size=Pt(54))
        add_body(slide, "Codebase Intelligence Report", y=Inches(3.2), color=ACCENT)
        add_body(slide, analysis_meta.get("repo_name", "Repository Analysis"), y=Inches(4.0), color=WHITE)
        add_body(slide, f"Generated: {datetime.now().strftime('%B %d, %Y')}", y=Inches(5.5))

        # Slide 2: Executive Summary
        slide = add_dark_slide("")
        add_title(slide, "Executive Summary")
        scores_text = (
            f"Scalability Score: {report_data.get('scalability_score', 'N/A')}/100\n"
            f"Production Readiness: {report_data.get('production_readiness_score', 'N/A')}/100\n"
            f"Security Findings: {len(report_data.get('security_findings', []))}\n"
            f"Files Analyzed: {analysis_meta.get('file_count', 0)}"
        )
        add_body(slide, scores_text, y=Inches(1.9), color=WHITE)

        # Slide 3: Architecture
        slide = add_dark_slide("")
        add_title(slide, "Architecture Overview")
        arch = report_data.get("architecture_summary", "")
        add_body(slide, arch[:600] + ("..." if len(arch) > 600 else ""), y=Inches(1.9))

        # Slide 4: Tech Stack
        slide = add_dark_slide("")
        add_title(slide, "Technology Stack")
        tech = report_data.get("tech_stack", [])
        add_body(slide, " · ".join(tech) if tech else "Not detected", y=Inches(1.9), color=ACCENT)

        # Slide 5: Security Findings
        slide = add_dark_slide("")
        add_title(slide, "Security Analysis")
        findings = report_data.get("security_findings", [])
        if findings:
            lines = []
            for f in findings[:8]:
                severity = f.get("severity", "?").upper()
                lines.append(f"[{severity}] {f.get('id', '?')} — {f.get('file', '?')}")
            add_body(slide, "\n".join(lines), y=Inches(1.9), color=WHITE)
        else:
            add_body(slide, "✓ No security issues detected", y=Inches(1.9), color=ACCENT)

        # Slides 6-8: Interview Questions
        questions = report_data.get("interview_questions", [])
        for i, qa in enumerate(questions[:6]):
            slide = add_dark_slide("")
            add_title(slide, f"Interview Q{i + 1}: {qa.get('category', '').title()}", size=Pt(22))
            add_body(slide, qa.get("question", ""), y=Inches(1.7), color=ACCENT)
            answer = qa.get("expected_answer", "")
            add_body(slide, answer[:400] + ("..." if len(answer) > 400 else ""), y=Inches(2.8))

        # Slide: Recommendations
        slide = add_dark_slide("")
        add_title(slide, "Recommendations")
        recs = report_data.get("recommendations", [])
        if recs:
            lines = [f"[{r.get('priority', '?').upper()}] {r.get('title', '?')}" for r in recs[:8]]
            add_body(slide, "\n".join(lines), y=Inches(1.9), color=WHITE)

        # Final: CTA
        slide = add_dark_slide("")
        add_title(slide, "ArchDefend", y=Inches(2.0), size=Pt(48))
        add_body(slide, "archdefend.io · Enterprise Codebase Intelligence", y=Inches(4.0), color=ACCENT)
        add_body(slide, "by EMEMZYVISUALS DIGITALS", y=Inches(5.0))

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        pptx_bytes = buf.read()
        logger.info(f"PPTX generated: {len(pptx_bytes):,} bytes")
        return pptx_bytes

    async def generate_markdown(self, report_data: dict, analysis_meta: dict) -> str:
        """Generate Markdown report."""
        lines = [
            f"# ArchDefend — Codebase Intelligence Report",
            f"",
            f"**Repository:** `{analysis_meta.get('repo_url', 'N/A')}`  ",
            f"**Generated:** {datetime.now().strftime('%Y-%m-%d %H:%M UTC')}  ",
            f"**By:** [EMEMZYVISUALS DIGITALS](https://archdefend.io)",
            f"",
            f"---",
            f"",
            f"## Executive Summary",
            f"",
            f"| Metric | Score |",
            f"|--------|-------|",
            f"| Scalability | {report_data.get('scalability_score', 'N/A')}/100 |",
            f"| Production Readiness | {report_data.get('production_readiness_score', 'N/A')}/100 |",
            f"| Security Findings | {len(report_data.get('security_findings', []))} |",
            f"| Files Analyzed | {analysis_meta.get('file_count', 0)} |",
            f"",
            f"## Architecture Summary",
            f"",
            report_data.get("architecture_summary", "Not available."),
            f"",
            f"## Technology Stack",
            f"",
            ", ".join(f"`{t}`" for t in report_data.get("tech_stack", [])),
            f"",
            f"## Security Findings",
            f"",
        ]

        for finding in report_data.get("security_findings", []):
            lines += [
                f"### [{finding.get('severity', '?').upper()}] {finding.get('id', '?')}",
                f"",
                f"**File:** `{finding.get('file', 'N/A')}`:{finding.get('line', '?')}  ",
                f"**CWE:** {finding.get('cwe', 'N/A')}  ",
                f"",
                finding.get("description", ""),
                f"",
                f"> **Fix:** {finding.get('remediation', 'Review and remediate.')}",
                f"",
            ]

        lines += [
            f"## API Inventory",
            f"",
            f"| Method | Path | Framework |",
            f"|--------|------|-----------|",
        ]
        for route in report_data.get("api_inventory", []):
            lines.append(f"| `{route.get('method', 'ANY')}` | `{route.get('path', '?')}` | {route.get('framework', '?')} |")

        lines += [f"", f"## Interview Defense Questions", f""]
        for i, qa in enumerate(report_data.get("interview_questions", []), 1):
            lines += [
                f"### Q{i}: {qa.get('question', '')}",
                f"",
                f"**Difficulty:** {qa.get('difficulty', '?')} · **Category:** {qa.get('category', '?')}",
                f"",
                f"**Answer:**",
                qa.get("expected_answer", ""),
                f"",
            ]

        lines += [f"---", f"", f"*Generated by [ArchDefend](https://archdefend.io) — EMEMZYVISUALS DIGITALS*"]
        return "\n".join(lines)

    def _build_pdf_html(self, report_data: dict, analysis_meta: dict) -> str:
        """Build HTML for PDF generation."""
        findings = report_data.get("security_findings", [])
        security_color = "#DC2626" if len(findings) > 5 else "#F59E0B" if findings else "#10B981"

        # Security HTML
        security_html = ""
        if findings:
            for f in findings[:15]:
                badge_class = f"badge-{f.get('severity', 'info')}"
                security_html += f"""
                <div class="finding">
                  <div class="finding-title">
                    <span class="badge {badge_class}">{f.get('severity', '?').upper()}</span>
                    {f.get('id', '?')} — {f.get('title', f.get('description', '')[:80])}
                  </div>
                  <div class="finding-meta">{f.get('file', 'N/A')} · {f.get('cwe', 'N/A')}</div>
                  <p>{f.get('description', '')}</p>
                  <p><strong>Fix:</strong> {f.get('remediation', '')}</p>
                </div>"""
        else:
            security_html = "<p style='color: #10B981; font-weight: 600;'>✓ No security issues detected by static analysis.</p>"

        # Tech stack HTML
        tech_html = " ".join(f'<span class="tech-pill">{t}</span>' for t in report_data.get("tech_stack", []))

        # API HTML
        routes = report_data.get("api_inventory", [])
        if routes:
            api_html = "<table><tr><th>Method</th><th>Path</th><th>Framework</th></tr>"
            for r in routes[:30]:
                api_html += f"<tr><td><code>{r.get('method', 'ANY')}</code></td><td><code>{r.get('path', '?')}</code></td><td>{r.get('framework', '?')}</td></tr>"
            api_html += "</table>"
        else:
            api_html = "<p>No API routes detected automatically. Review source files for route definitions.</p>"

        # Interview HTML
        interview_html = ""
        for i, qa in enumerate(report_data.get("interview_questions", []), 1):
            interview_html += f"""
            <div class="qa-item">
              <div class="qa-question">Q{i}: {qa.get('question', '')}</div>
              <div class="qa-answer">{qa.get('expected_answer', '')}</div>
            </div>"""

        # Recommendations HTML
        recs = report_data.get("recommendations", [])
        rec_rows = ""
        for r in recs:
            priority_color = {"critical": "#DC2626", "high": "#F59E0B", "medium": "#CA8A04", "low": "#059669"}.get(r.get("priority", ""), "#6B7280")
            rec_rows += f"""
            <tr>
              <td><strong style='color:{priority_color}'>{r.get('priority', '?').upper()}</strong></td>
              <td>{r.get('title', '')}</td>
              <td>{r.get('category', '')}</td>
              <td>{r.get('effort', '?')}</td>
            </tr>"""
        recommendations_html = f"""
        <table>
          <tr><th>Priority</th><th>Title</th><th>Category</th><th>Effort</th></tr>
          {rec_rows}
        </table>""" if recs else "<p>No recommendations generated.</p>"

        return PDF_HTML_TEMPLATE.format(
            repo_name=analysis_meta.get("repo_name", "Repository"),
            repo_url=analysis_meta.get("repo_url", ""),
            generated_at=datetime.now().strftime("%Y-%m-%d %H:%M UTC"),
            scalability_score=report_data.get("scalability_score", "—"),
            readiness_score=report_data.get("production_readiness_score", "—"),
            security_count=len(findings),
            security_color=security_color,
            file_count=analysis_meta.get("file_count", 0),
            architecture_summary=report_data.get("architecture_summary", "").replace("\n", "<br>"),
            tech_stack_html=tech_html,
            security_html=security_html,
            api_html=api_html,
            interview_html=interview_html,
            recommendations_html=recommendations_html,
        )


# Singleton
export_engine = ExportEngine()
