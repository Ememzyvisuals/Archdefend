"""
ArchDefend — Export Engine
Generates PDF, PPTX, Markdown, and HTML from analysis reports.
"""
import io
import logging
from datetime import datetime

logger = logging.getLogger("archdefend.export")

_PDF_TEMPLATE = """<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700&family=JetBrains+Mono:wght@400;500&display=swap');
*{{box-sizing:border-box;margin:0;padding:0}}
body{{font-family:'Inter',sans-serif;background:#fff;color:#111;font-size:11pt;line-height:1.6}}
.cover{{background:linear-gradient(135deg,#07070A,#111827);color:#F0F0F5;padding:60px 48px;min-height:100vh;display:flex;flex-direction:column;justify-content:center}}
.cover h1{{font-size:40pt;font-weight:700;margin-bottom:12px;letter-spacing:-1px}}
.cover h1 span{{color:#22D3EE}}
.cover p{{color:#5A5A72;font-size:13pt;margin-bottom:6px}}
.cover .repo{{color:#22D3EE;font-size:18pt;font-weight:600;margin-top:16px}}
.cover .meta{{margin-top:48px;padding-top:24px;border-top:1px solid #1E1E2A;color:#404058;font-size:9pt;font-family:'JetBrains Mono',monospace}}
.page{{padding:48px;page-break-after:always}}
h2{{font-size:22pt;font-weight:700;margin-bottom:24px;padding-bottom:12px;border-bottom:2px solid #22D3EE;color:#111}}
h3{{font-size:14pt;font-weight:600;margin:24px 0 12px;color:#111}}
p{{margin-bottom:12px;color:#374151}}
.scores{{display:flex;gap:16px;margin:24px 0;flex-wrap:wrap}}
.score-box{{border:2px solid #E5E7EB;border-radius:12px;padding:20px 28px;text-align:center;min-width:140px}}
.score-num{{font-size:36pt;font-weight:700;color:#22D3EE}}
.score-lbl{{font-size:9pt;color:#6B7280;margin-top:4px}}
.badge{{display:inline-block;padding:3px 10px;border-radius:20px;font-size:8pt;font-weight:600;margin:2px}}
.badge-critical{{background:#FEE2E2;color:#DC2626}}
.badge-high{{background:#FEF3C7;color:#D97706}}
.badge-medium{{background:#FEF9C3;color:#CA8A04}}
.badge-low{{background:#D1FAE5;color:#059669}}
.finding{{border:1px solid #E5E7EB;border-radius:8px;padding:16px;margin-bottom:12px;background:#F9FAFB}}
.finding-title{{font-weight:600;font-size:10pt;margin-bottom:6px}}
.finding-meta{{font-size:8pt;color:#6B7280;font-family:'JetBrains Mono',monospace;margin-bottom:8px}}
.qa{{border:1px solid #E5E7EB;border-radius:8px;padding:20px;margin-bottom:16px}}
.qa-q{{font-weight:600;font-size:10pt;margin-bottom:10px;color:#111}}
.qa-a{{color:#374151;font-size:9pt;line-height:1.75}}
.tech-pill{{display:inline-block;padding:4px 12px;border-radius:20px;border:1px solid #E5E7EB;font-size:8pt;color:#374151;margin:3px;background:#F9FAFB}}
table{{width:100%;border-collapse:collapse;margin:16px 0}}
th{{background:#F9FAFB;padding:10px 12px;text-align:left;font-size:8pt;font-weight:600;color:#6B7280;border-bottom:1px solid #E5E7EB;text-transform:uppercase;letter-spacing:.5px}}
td{{padding:10px 12px;font-size:9pt;border-bottom:1px solid #F3F4F6;color:#374151}}
.footer{{text-align:center;color:#9CA3AF;font-size:8pt;padding:24px;border-top:1px solid #E5E7EB;font-family:'JetBrains Mono',monospace}}
</style>
</head>
<body>
<div class="cover">
  <h1>Arch<span>Defend</span></h1>
  <p>Codebase Intelligence Report</p>
  <p class="repo">{repo_name}</p>
  <p>{repo_url}</p>
  <div class="meta">Generated: {generated_at}<br>ArchDefend · archdefend.io · EMEMZYVISUALS DIGITALS</div>
</div>
<div class="page">
  <h2>Executive Summary</h2>
  <div class="scores">
    <div class="score-box"><div class="score-num">{scalability}</div><div class="score-lbl">Scalability</div></div>
    <div class="score-box"><div class="score-num">{readiness}</div><div class="score-lbl">Production Ready</div></div>
    <div class="score-box"><div class="score-num" style="color:{'#DC2626' if int(str(security_count).replace('—','0') or 0) > 5 else '#F59E0B' if int(str(security_count).replace('—','0') or 0) > 0 else '#10B981'}">{security_count}</div><div class="score-lbl">Security Findings</div></div>
    <div class="score-box"><div class="score-num">{file_count}</div><div class="score-lbl">Files Analyzed</div></div>
  </div>
  <h3>Architecture Summary</h3>
  <p>{arch_summary}</p>
  <h3>Technology Stack</h3>
  <div>{tech_html}</div>
</div>
<div class="page">
  <h2>Security Analysis</h2>
  {security_html}
</div>
<div class="page">
  <h2>API Inventory</h2>
  {api_html}
</div>
<div class="page">
  <h2>Interview Defense Questions</h2>
  {interview_html}
</div>
<div class="page">
  <h2>Recommendations</h2>
  {rec_html}
  <div class="footer">ArchDefend · archdefend.io · EMEMZYVISUALS DIGITALS</div>
</div>
</body>
</html>"""


class ExportEngine:
    async def pdf(self, report: dict, meta: dict) -> bytes:
        try:
            from weasyprint import HTML
        except ImportError:
            raise RuntimeError("weasyprint not installed")
        html = self._pdf_html(report, meta)
        return HTML(string=html).write_pdf()

    async def pptx(self, report: dict, meta: dict) -> bytes:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
        except ImportError:
            raise RuntimeError("python-pptx not installed")

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        BG = RGBColor(7, 7, 10)
        CYAN = RGBColor(34, 211, 238)
        WHITE = RGBColor(238, 238, 245)
        MUTED = RGBColor(90, 90, 114)

        def slide(title: str, body: str, body_color=None):
            sl = prs.slides.add_slide(prs.slide_layouts[6])
            sl.background.fill.solid()
            sl.background.fill.fore_color.rgb = BG
            for pos, text, size, color in [
                (0.6, title, Pt(28), WHITE),
                (0.6, body, Pt(13), body_color or MUTED),
            ]:
                tb = sl.shapes.add_textbox(Inches(0.7), Inches(pos if text == title else 1.8), Inches(11.8), Inches(5))
                tf = tb.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = text[:700]
                run.font.size = size
                run.font.color.rgb = color
                run.font.bold = (text == title)
            return sl

        slide("ArchDefend — Codebase Report", meta.get("repo_name", ""), CYAN)
        slide("Executive Summary",
              f"Scalability: {report.get('scalability_score','—')}/100\n"
              f"Readiness: {report.get('production_readiness_score','—')}/100\n"
              f"Security: {len(report.get('security_findings',[]))} findings\n"
              f"Files: {meta.get('file_count',0)}", WHITE)
        slide("Architecture Overview", report.get("architecture_summary","")[:700])
        slide("Technology Stack", " · ".join(report.get("tech_stack", [])), CYAN)

        for i, finding in enumerate(report.get("security_findings", [])[:8]):
            slide(f"[{finding.get('severity','?').upper()}] {finding.get('id','?')}",
                  f"{finding.get('file','?')}\n{finding.get('description','')}\nFix: {finding.get('remediation','')}")

        for i, qa in enumerate(report.get("interview_questions", [])[:6]):
            slide(f"Q{i+1}: {qa.get('category','?').title()}", qa.get("question",""))
            slide("Expected Answer", qa.get("expected_answer","")[:700])

        slide("Recommendations",
              "\n".join(f"[{r.get('priority','?').upper()}] {r.get('title','?')}" for r in report.get("recommendations",[])[:8]))
        slide("ArchDefend", "archdefend.io · EMEMZYVISUALS DIGITALS", CYAN)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.read()

    async def markdown(self, report: dict, meta: dict) -> str:
        lines = [
            "# ArchDefend — Codebase Intelligence Report",
            f"\n**Repository:** `{meta.get('repo_url','')}`  ",
            f"**Generated:** {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}  ",
            f"**By:** [EMEMZYVISUALS DIGITALS](https://archdefend.io)\n\n---\n",
            "## Executive Summary\n",
            "| Metric | Score |","|---|---|",
            f"| Scalability | {report.get('scalability_score','—')}/100 |",
            f"| Production Readiness | {report.get('production_readiness_score','—')}/100 |",
            f"| Security Findings | {len(report.get('security_findings',[]))} |",
            f"| Files Analyzed | {meta.get('file_count',0)} |",
            f"\n## Architecture Summary\n\n{report.get('architecture_summary','Not available.')}\n",
            f"\n## Technology Stack\n\n{', '.join(f'`{t}`' for t in report.get('tech_stack',[]))}\n",
            "\n## Security Findings\n",
        ]
        for f in report.get("security_findings", []):
            lines += [
                f"\n### [{f.get('severity','?').upper()}] {f.get('id','?')}",
                f"**File:** `{f.get('file','?')}`:{f.get('line','?')} · **CWE:** {f.get('cwe','?')}",
                f"\n{f.get('description','')}\n",
                f"> **Remediation:** {f.get('remediation','Review and fix.')}\n",
            ]
        lines += ["\n## API Inventory\n\n| Method | Path | Framework |","|---|---|---|"]
        for r in report.get("api_inventory", []):
            lines.append(f"| `{r.get('method','?')}` | `{r.get('path','?')}` | {r.get('framework','?')} |")
        lines.append("\n## Interview Defense Questions\n")
        for i, qa in enumerate(report.get("interview_questions", []), 1):
            lines += [
                f"\n### Q{i}: {qa.get('question','')}",
                f"**Category:** {qa.get('category','?')} · **Difficulty:** {qa.get('difficulty','?')}\n",
                f"**Answer:**\n{qa.get('expected_answer','')}\n",
            ]
        lines.append("\n---\n\n*Generated by [ArchDefend](https://archdefend.io) — EMEMZYVISUALS DIGITALS*")
        return "\n".join(lines)

    def _pdf_html(self, r: dict, m: dict) -> str:
        findings = r.get("security_findings", [])
        sc = len(findings)
        sc_color = "#DC2626" if sc > 5 else "#F59E0B" if sc > 0 else "#10B981"

        sec_html = ""
        for f in findings[:15]:
            sec_html += f'<div class="finding"><div class="finding-title"><span class="badge badge-{f.get("severity","info")}">{f.get("severity","?").upper()}</span> {f.get("id","?")} — {f.get("title",f.get("description",""))[:80]}</div><div class="finding-meta">{f.get("file","?")} · {f.get("cwe","?")}</div><p>{f.get("description","")}</p><p><strong>Fix:</strong> {f.get("remediation","")}</p></div>'
        if not sec_html:
            sec_html = '<p style="color:#10B981;font-weight:600">✓ No security issues detected.</p>'

        routes = r.get("api_inventory", [])
        api_html = '<table><tr><th>Method</th><th>Path</th><th>Framework</th></tr>' + "".join(f'<tr><td><code>{x.get("method","?")}</code></td><td><code>{x.get("path","?")}</code></td><td>{x.get("framework","?")}</td></tr>' for x in routes[:30]) + "</table>" if routes else "<p>No API routes detected automatically.</p>"

        interview_html = "".join(
            f'<div class="qa"><div class="qa-q">Q{i}: {qa.get("question","")}</div><div class="qa-a">{qa.get("expected_answer","")}</div></div>'
            for i, qa in enumerate(r.get("interview_questions", []), 1)
        )

        recs = r.get("recommendations", [])
        _pc = {"critical": "#DC2626", "high": "#D97706", "medium": "#CA8A04", "low": "#059669"}
        def _rec_row(x):
            c = _pc.get(x.get("priority", ""), "#6B7280")
            p = x.get("priority", "?").upper()
            t = x.get("title", "")
            cat = x.get("category", "")
            e = x.get("effort", "?")
            return f'<tr><td><strong style="color:{c}">{p}</strong></td><td>{t}</td><td>{cat}</td><td>{e}</td></tr>'
        rec_html = '<table><tr><th>Priority</th><th>Title</th><th>Category</th><th>Effort</th></tr>' + "".join(
            _rec_row(x) for x in recs
        ) + "</table>" if recs else "<p>No recommendations generated.</p>" 

        tech_html = "".join(f'<span class="tech-pill">{t}</span>' for t in r.get("tech_stack", []))

        return _PDF_TEMPLATE.format(
            repo_name=m.get("repo_name", "Repository"),
            repo_url=m.get("repo_url", ""),
            generated_at=datetime.utcnow().strftime("%Y-%m-%d %H:%M UTC"),
            scalability=r.get("scalability_score", "—"),
            readiness=r.get("production_readiness_score", "—"),
            security_count=sc,
            file_count=m.get("file_count", 0),
            arch_summary=r.get("architecture_summary", "").replace("\n", "<br>"),
            tech_html=tech_html,
            security_html=sec_html,
            api_html=api_html,
            interview_html=interview_html,
            rec_html=rec_html,
        )


export_engine = ExportEngine()
